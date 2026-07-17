"""取引事例レポートを PowerPoint / Word 形式で生成する。

PDF と異なりフォントはクライアント（PowerPoint / Word）側で解決されるため、
日本語フォントの埋め込みは不要。仲介業者が客先で編集・追記できる点も狙い。
種別（seller / buyer / appraisal）ごとに `ReportContext.sections` で指定された
セクションを、その順序で描画する。
"""

from __future__ import annotations

from datetime import datetime
from io import BytesIO
from typing import Optional

from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml.ns import qn as docx_qn
from docx.shared import Pt as DocxPt
from docx.shared import RGBColor as DocxColor
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn as pptx_qn
from pptx.util import Emu, Inches
from pptx.util import Pt as PptPt

from app.api.schemas import MunicipalityDetail, ReportContext
from app.web.formatters import (
    format_count,
    format_man_yen,
    format_percent,
    format_yen_per_sqm,
    quarter_label,
)

FONT = "Meiryo"
PRIMARY = RGBColor(0x03, 0x69, 0xA1)
DARK = RGBColor(0x0F, 0x17, 0x2A)
GREY = RGBColor(0x64, 0x74, 0x8B)
LIGHT = RGBColor(0xF1, 0xF5, 0xF9)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DOCX_GREY = DocxColor(0x64, 0x74, 0x8B)

SOURCE_NOTE = "出典: 国土交通省 不動産情報ライブラリ（不動産取引価格情報・地価公示）"
DISCLAIMER = "本レポートは参考情報です。実際の取引条件は個別の物件・取引により異なります。"

METHODOLOGY_TEXT = (
    "本資料の査定価格は、国土交通省 不動産情報ライブラリに登録された実際の取引価格情報"
    "および地価公示データに基づく参考値です。物件種別ごとの㎡単価・成約事例・地価水準を"
    "総合的に勘案していますが、個別要因（築年数・方位・接道・室内状態・売却時期等）により"
    "実際の成約価格は変動します。正式な査定は現地調査のうえ行ってください。"
)


def _recent_row(tx) -> list[str]:
    period = tx.period_label or quarter_label(tx.trade_year, tx.trade_quarter)
    area = f"{tx.area:,.0f}㎡" if tx.area else "—"
    return [
        period,
        tx.property_type or "—",
        tx.district_name or "—",
        area,
        format_man_yen(tx.trade_price),
        format_yen_per_sqm(tx.unit_price),
    ]


RECENT_HEADERS = ["時期", "種別", "地区", "面積", "取引価格", "㎡単価"]
PROPERTY_HEADERS = ["種別", "件数", "平均価格", "㎡単価"]
BRACKET_HEADERS = ["価格帯", "件数", "平均㎡単価"]
YEARLY_HEADERS = ["年", "件数", "平均取引価格"]
LAND_TREND_HEADERS = ["調査年", "地点数", "平均地価", "前年比"]


def _section_title(section: str) -> str:
    return {
        "summary": "エリアサマリー",
        "recent_cases": "直近の取引事例",
        "price_brackets": "価格帯別の分布",
        "property_mix": "物件種別の内訳（最新四半期）",
        "yearly_trend": "取引価格の年次推移",
        "land_price_trend": "地価公示の推移",
        "land_price": "地価公示サマリー",
        "methodology": "査定の手法と留意事項",
    }.get(section, section)


# --------------------------------------------------------------------------- #
# PowerPoint
# --------------------------------------------------------------------------- #


def _style_run(run, size: float, bold: bool = False, color: Optional[RGBColor] = None) -> None:
    font = run.font
    font.size = PptPt(size)
    font.bold = bold
    font.name = FONT
    if color is not None:
        font.color.rgb = color
    rpr = font._rPr
    if rpr is not None:
        for tag in ("a:ea", "a:cs"):
            el = rpr.find(pptx_qn(tag))
            if el is None:
                el = rpr.makeelement(pptx_qn(tag), {"typeface": FONT})
                rpr.append(el)
            else:
                el.set("typeface", FONT)


def _add_text(
    slide,
    left,
    top,
    width,
    height,
    text: str,
    size: float,
    *,
    bold: bool = False,
    color: Optional[RGBColor] = None,
    align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    para = tf.paragraphs[0]
    para.alignment = align
    run = para.add_run()
    run.text = text
    _style_run(run, size, bold=bold, color=color)
    return box


def _set_cell(cell, text: str, *, size: float, bold=False, color=None, fill=None, align=PP_ALIGN.LEFT):
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    cell.margin_left = Emu(64008)
    cell.margin_right = Emu(64008)
    cell.margin_top = Emu(18288)
    cell.margin_bottom = Emu(18288)
    if fill is not None:
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill
    tf = cell.text_frame
    tf.word_wrap = True
    para = tf.paragraphs[0]
    para.alignment = align
    run = para.add_run()
    run.text = text
    _style_run(run, size, bold=bold, color=color or DARK)


def _add_table(slide, left, top, width, headers: list[str], rows: list[list[str]], right_align_from: int = 1):
    n_rows = len(rows) + 1
    n_cols = len(headers)
    height = Inches(0.34) * n_rows
    graphic = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    table = graphic.table
    table.first_row = True
    for c, head in enumerate(headers):
        align = PP_ALIGN.LEFT if c < right_align_from else PP_ALIGN.RIGHT
        _set_cell(table.cell(0, c), head, size=11, bold=True, color=WHITE, fill=PRIMARY, align=align)
    for r, row in enumerate(rows, start=1):
        fill = WHITE if r % 2 else LIGHT
        for c, value in enumerate(row):
            align = PP_ALIGN.LEFT if c < right_align_from else PP_ALIGN.RIGHT
            _set_cell(table.cell(r, c), value, size=10.5, fill=fill, align=align)
    return graphic


def _add_line_chart(slide, left, top, width, height, categories, series_name, values, unit=""):
    data = CategoryChartData()
    data.categories = [str(c) for c in categories]
    data.add_series(series_name, values)
    graphic = slide.shapes.add_chart(XL_CHART_TYPE.LINE_MARKERS, left, top, width, height, data)
    chart = graphic.chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    chart.font.name = FONT
    chart.font.size = PptPt(10)
    return chart


def _content_slide(prs, layout, title: str):
    slide = prs.slides.add_slide(layout)
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.14))
    band.fill.solid()
    band.fill.fore_color.rgb = PRIMARY
    band.line.fill.background()
    band.shadow.inherit = False
    _add_text(slide, Inches(0.6), Inches(0.4), prs.slide_width - Inches(1.2), Inches(0.7),
              title, 22, bold=True, color=DARK)
    return slide


def _pptx_cover(prs, layout, detail: MunicipalityDetail, report: ReportContext) -> None:
    slide = prs.slides.add_slide(layout)
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(3.4))
    band.fill.solid()
    band.fill.fore_color.rgb = PRIMARY
    band.line.fill.background()
    band.shadow.inherit = False

    _add_text(slide, Inches(0.9), Inches(0.7), prs.slide_width - Inches(1.8), Inches(0.5),
              report.report_type_label, 14, color=RGBColor(0xBA, 0xE6, 0xFD))
    _add_text(slide, Inches(0.9), Inches(1.2), prs.slide_width - Inches(1.8), Inches(1.0),
              f"{detail.prefecture_name}{detail.name_ja}", 40, bold=True, color=WHITE)
    _add_text(slide, Inches(0.9), Inches(2.35), prs.slide_width - Inches(1.8), Inches(0.6),
              f"周辺取引価格の統計・事例に基づくエリア分析（{report.period_label}）",
              16, color=RGBColor(0xE0, 0xF2, 0xFE))

    kpis = [
        ("累計取引件数", f"{format_count(detail.total_transactions)}件"),
        ("平均取引価格", format_man_yen(detail.recent_avg_price)),
    ]
    if detail.latest_year:
        kpis.append(("最新データ", quarter_label(detail.latest_year, detail.latest_quarter or 1)))
    box_w = Inches(3.6)
    for i, (label, value) in enumerate(kpis):
        left = Inches(0.9) + box_w * i + Inches(0.25) * i
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(3.9), box_w, Inches(1.3))
        card.fill.solid()
        card.fill.fore_color.rgb = LIGHT
        card.line.color.rgb = RGBColor(0xE2, 0xE8, 0xF0)
        card.shadow.inherit = False
        _add_text(slide, left + Inches(0.25), Inches(4.05), box_w - Inches(0.5), Inches(0.4),
                  label, 12, color=GREY)
        _add_text(slide, left + Inches(0.25), Inches(4.45), box_w - Inches(0.5), Inches(0.6),
                  value, 22, bold=True, color=PRIMARY)

    footer_bits = [SOURCE_NOTE, f"作成日: {datetime.now().strftime('%Y年%m月%d日')}"]
    if detail.stats_updated_at:
        footer_bits.append(f"データ更新: {detail.stats_updated_at.strftime('%Y年%m月%d日')}")
    _add_text(slide, Inches(0.9), Inches(6.7), prs.slide_width - Inches(1.8), Inches(0.5),
              "　｜　".join(footer_bits), 10, color=GREY)


def _pptx_section(prs, layout, detail: MunicipalityDetail, report: ReportContext, section: str) -> None:
    left = Inches(0.6)
    top = Inches(1.35)
    width = prs.slide_width - Inches(1.2)

    if section == "summary":
        slide = _content_slide(prs, layout, _section_title(section))
        _add_text(slide, left, top, width, Inches(4.5), report.summary_text, 15, color=DARK)
        return

    if section == "recent_cases":
        if not report.recent_transactions:
            return
        slide = _content_slide(prs, layout, _section_title(section))
        shown = report.recent_transactions[:14]
        rows = [_recent_row(tx) for tx in shown]
        _add_table(slide, left, top, width, RECENT_HEADERS, rows, right_align_from=3)
        extra = len(report.recent_transactions) - len(shown)
        if extra > 0:
            _add_text(slide, left, Inches(7.0), width, Inches(0.35),
                      f"ほか {extra} 件の取引事例があります。", 10, color=GREY)
        return

    if section == "price_brackets":
        if not report.price_brackets:
            return
        slide = _content_slide(prs, layout, _section_title(section))
        rows = [
            [b.label, f"{format_count(b.transaction_count)}件", format_yen_per_sqm(b.unit_price_avg)]
            for b in report.price_brackets
        ]
        _add_table(slide, left, top, Inches(7.5), BRACKET_HEADERS, rows)
        return

    if section == "property_mix":
        if not report.property_stats:
            return
        slide = _content_slide(prs, layout, _section_title(section))
        rows = [
            [
                s.property_type or "—",
                f"{format_count(s.transaction_count)}件",
                format_man_yen(s.trade_price_avg),
                format_yen_per_sqm(s.unit_price_avg),
            ]
            for s in report.property_stats
        ]
        _add_table(slide, left, top, width, PROPERTY_HEADERS, rows)
        return

    if section == "yearly_trend":
        stats = report.yearly_stats[-10:]
        if not stats:
            return
        slide = _content_slide(prs, layout, _section_title(section))
        years = [f"{s.trade_year}年" for s in stats]
        values = [round(s.trade_price_avg / 10_000) if s.trade_price_avg else None for s in stats]
        _add_line_chart(slide, left, top, Inches(7.2), Inches(4.6), years, "平均取引価格（万円）", values)
        rows = [
            [f"{s.trade_year}年", f"{format_count(s.transaction_count)}件", format_man_yen(s.trade_price_avg)]
            for s in stats
        ]
        _add_table(slide, Inches(8.1), top, Inches(4.6), YEARLY_HEADERS, rows)
        return

    if section == "land_price_trend":
        stats = report.land_price_yearly[-10:]
        if not stats:
            return
        slide = _content_slide(prs, layout, _section_title(section))
        years = [f"{s.survey_year}年" for s in stats]
        values = [round(s.avg_unit_price) if s.avg_unit_price else None for s in stats]
        _add_line_chart(slide, left, top, Inches(7.2), Inches(4.6), years, "平均地価（円/㎡）", values)
        rows = [
            [
                f"{s.survey_year}年",
                f"{format_count(s.point_count)}地点",
                format_yen_per_sqm(s.avg_unit_price),
                format_percent(s.yoy_avg_price_pct),
            ]
            for s in stats
        ]
        _add_table(slide, Inches(8.1), top, Inches(4.6), LAND_TREND_HEADERS, rows, right_align_from=1)
        return

    if section == "land_price":
        land = detail.land_prices
        if not land or not land.point_count:
            return
        slide = _content_slide(prs, layout, _section_title(section))
        cards = [
            ("地点数", f"{format_count(land.point_count)}地点"),
            ("平均地価", format_yen_per_sqm(land.avg_unit_price)),
            ("前年比", format_percent(land.yoy_change_avg)),
        ]
        box_w = Inches(3.8)
        for i, (label, value) in enumerate(cards):
            cleft = left + box_w * i + Inches(0.25) * i
            card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, cleft, top, box_w, Inches(1.5))
            card.fill.solid()
            card.fill.fore_color.rgb = LIGHT
            card.line.color.rgb = RGBColor(0xE2, 0xE8, 0xF0)
            card.shadow.inherit = False
            _add_text(slide, cleft + Inches(0.3), top + Inches(0.2), box_w - Inches(0.6), Inches(0.4),
                      label, 13, color=GREY)
            _add_text(slide, cleft + Inches(0.3), top + Inches(0.65), box_w - Inches(0.6), Inches(0.7),
                      value, 24, bold=True, color=PRIMARY)
        if land.latest_year:
            _add_text(slide, left, top + Inches(1.8), width, Inches(0.4),
                      f"{land.latest_year}年 地価公示に基づく。", 11, color=GREY)
        return

    if section == "methodology":
        slide = _content_slide(prs, layout, _section_title(section))
        _add_text(slide, left, top, width, Inches(4.5), METHODOLOGY_TEXT, 14, color=DARK)
        _add_text(slide, left, Inches(6.4), width, Inches(0.8), f"{SOURCE_NOTE}\n{DISCLAIMER}",
                  10, color=GREY)
        return


def build_pptx(detail: MunicipalityDetail, report: ReportContext) -> bytes:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    layout = prs.slide_layouts[6]  # blank

    _pptx_cover(prs, layout, detail, report)
    for section in report.sections:
        _pptx_section(prs, layout, detail, report, section)

    buffer = BytesIO()
    prs.save(buffer)
    return buffer.getvalue()


# --------------------------------------------------------------------------- #
# Word
# --------------------------------------------------------------------------- #


def _set_style_font(style, size: Optional[float] = None, color: Optional[DocxColor] = None) -> None:
    style.font.name = FONT
    if size is not None:
        style.font.size = DocxPt(size)
    if color is not None:
        style.font.color.rgb = color
    rpr = style.element.get_or_add_rPr()
    rfonts = rpr.get_or_add_rFonts()
    rfonts.set(docx_qn("w:eastAsia"), FONT)


def _docx_table(doc, headers: list[str], rows: list[list[str]], right_align_from: int = 1):
    table = doc.add_table(rows=1, cols=len(headers))
    table.style = "Light Grid Accent 1"
    for c, head in enumerate(headers):
        cell = table.rows[0].cells[c]
        cell.text = head
        para = cell.paragraphs[0]
        para.alignment = WD_ALIGN_PARAGRAPH.LEFT if c < right_align_from else WD_ALIGN_PARAGRAPH.RIGHT
        run = para.runs[0]
        run.font.bold = True
        run.font.name = FONT
        run.font.size = DocxPt(10)
        run._element.rPr.rFonts.set(docx_qn("w:eastAsia"), FONT)
    for row in rows:
        cells = table.add_row().cells
        for c, value in enumerate(row):
            cells[c].text = value
            para = cells[c].paragraphs[0]
            para.alignment = WD_ALIGN_PARAGRAPH.LEFT if c < right_align_from else WD_ALIGN_PARAGRAPH.RIGHT
            run = para.runs[0]
            run.font.name = FONT
            run.font.size = DocxPt(10)
            run._element.rPr.rFonts.set(docx_qn("w:eastAsia"), FONT)
    return table


def _docx_section(doc, detail: MunicipalityDetail, report: ReportContext, section: str) -> None:
    if section == "summary":
        doc.add_heading(_section_title(section), level=1)
        doc.add_paragraph(report.summary_text)
        return

    if section == "recent_cases":
        if not report.recent_transactions:
            return
        doc.add_heading(_section_title(section), level=1)
        rows = [_recent_row(tx) for tx in report.recent_transactions]
        _docx_table(doc, RECENT_HEADERS, rows, right_align_from=3)
        return

    if section == "price_brackets":
        if not report.price_brackets:
            return
        doc.add_heading(_section_title(section), level=1)
        rows = [
            [b.label, f"{format_count(b.transaction_count)}件", format_yen_per_sqm(b.unit_price_avg)]
            for b in report.price_brackets
        ]
        _docx_table(doc, BRACKET_HEADERS, rows)
        return

    if section == "property_mix":
        if not report.property_stats:
            return
        doc.add_heading(_section_title(section), level=1)
        rows = [
            [
                s.property_type or "—",
                f"{format_count(s.transaction_count)}件",
                format_man_yen(s.trade_price_avg),
                format_yen_per_sqm(s.unit_price_avg),
            ]
            for s in report.property_stats
        ]
        _docx_table(doc, PROPERTY_HEADERS, rows)
        return

    if section == "yearly_trend":
        stats = report.yearly_stats[-10:]
        if not stats:
            return
        doc.add_heading(_section_title(section), level=1)
        rows = [
            [f"{s.trade_year}年", f"{format_count(s.transaction_count)}件", format_man_yen(s.trade_price_avg)]
            for s in stats
        ]
        _docx_table(doc, YEARLY_HEADERS, rows)
        return

    if section == "land_price_trend":
        stats = report.land_price_yearly[-10:]
        if not stats:
            return
        doc.add_heading(_section_title(section), level=1)
        rows = [
            [
                f"{s.survey_year}年",
                f"{format_count(s.point_count)}地点",
                format_yen_per_sqm(s.avg_unit_price),
                format_percent(s.yoy_avg_price_pct),
            ]
            for s in stats
        ]
        _docx_table(doc, LAND_TREND_HEADERS, rows, right_align_from=1)
        return

    if section == "land_price":
        land = detail.land_prices
        if not land or not land.point_count:
            return
        doc.add_heading(_section_title(section), level=1)
        year_txt = f"（{land.latest_year}年）" if land.latest_year else ""
        doc.add_paragraph(
            f"地点数 {format_count(land.point_count)}地点{year_txt}　｜　"
            f"平均地価 {format_yen_per_sqm(land.avg_unit_price)}　｜　"
            f"前年比 {format_percent(land.yoy_change_avg)}"
        )
        return

    if section == "methodology":
        doc.add_heading(_section_title(section), level=1)
        doc.add_paragraph(METHODOLOGY_TEXT)
        return


def build_docx(detail: MunicipalityDetail, report: ReportContext) -> bytes:
    doc = Document()
    _set_style_font(doc.styles["Normal"], size=10.5)
    for style_name in ("Title", "Heading 1", "Heading 2"):
        try:
            _set_style_font(doc.styles[style_name])
        except KeyError:
            pass

    doc.add_heading(f"{detail.prefecture_name}{detail.name_ja}", level=0)
    subtitle = doc.add_paragraph()
    run = subtitle.add_run(f"{report.report_type_label}　｜　{report.period_label}")
    run.font.color.rgb = DOCX_GREY
    run.font.name = FONT
    run._element.rPr.rFonts.set(docx_qn("w:eastAsia"), FONT)

    kpi = doc.add_paragraph()
    kpi_parts = [
        f"累計取引件数: {format_count(detail.total_transactions)}件",
        f"平均取引価格: {format_man_yen(detail.recent_avg_price)}",
    ]
    if detail.latest_year:
        kpi_parts.append(f"最新データ: {quarter_label(detail.latest_year, detail.latest_quarter or 1)}")
    krun = kpi.add_run("　｜　".join(kpi_parts))
    krun.font.bold = True
    krun.font.name = FONT
    krun._element.rPr.rFonts.set(docx_qn("w:eastAsia"), FONT)

    for section in report.sections:
        _docx_section(doc, detail, report, section)

    doc.add_paragraph()
    footer_bits = [SOURCE_NOTE, DISCLAIMER, f"作成日: {datetime.now().strftime('%Y年%m月%d日')}"]
    if detail.stats_updated_at:
        footer_bits.append(f"データ更新: {detail.stats_updated_at.strftime('%Y年%m月%d日')}")
    for bit in footer_bits:
        para = doc.add_paragraph()
        run = para.add_run(bit)
        run.font.size = DocxPt(9)
        run.font.color.rgb = DOCX_GREY
        run.font.name = FONT
        run._element.rPr.rFonts.set(docx_qn("w:eastAsia"), FONT)

    buffer = BytesIO()
    doc.save(buffer)
    return buffer.getvalue()
